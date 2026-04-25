#!/usr/bin/env python3
"""slide2pdf — convert presentations (PPTX / ODP) to readable PDFs.

Usage
-----
    python slide2pdf.py INPUT [OUTPUT] [--libreoffice]

* INPUT   – path to .pptx or .odp file
* OUTPUT  – optional output path (default: same directory as INPUT, .pdf)
* --libreoffice – use LibreOffice headless as conversion backend (when available)

Without --libreoffice the tool uses a pure-Python pipeline:
  python-pptx  →  parse slide content
  reportlab    →  render readable PDF (text-first, high-fidelity layout)
"""

from __future__ import annotations

import argparse
import os
import subprocess
import sys
import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm, mm
from reportlab.platypus import (
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

_SLIDE_W, _SLIDE_H = landscape(A4)

_STYLES = getSampleStyleSheet()

_STYLE_TITLE = ParagraphStyle(
    "SlideTitle",
    parent=_STYLES["Heading1"],
    fontSize=22,
    leading=28,
    spaceAfter=8,
    textColor=colors.HexColor("#1F3864"),
)

_STYLE_SUBTITLE = ParagraphStyle(
    "SlideSubtitle",
    parent=_STYLES["Heading2"],
    fontSize=16,
    leading=20,
    spaceAfter=6,
    textColor=colors.HexColor("#2E75B6"),
)

_STYLE_BODY = ParagraphStyle(
    "SlideBody",
    parent=_STYLES["BodyText"],
    fontSize=11,
    leading=15,
    spaceAfter=4,
    wordWrap="LTR",
)

_STYLE_BULLET = ParagraphStyle(
    "SlideBullet",
    parent=_STYLE_BODY,
    leftIndent=12,
    bulletIndent=0,
    bulletFontName="Helvetica",
    bulletFontSize=11,
)

_STYLE_SLIDE_NUM = ParagraphStyle(
    "SlideNum",
    parent=_STYLES["Normal"],
    fontSize=8,
    textColor=colors.grey,
    alignment=2,  # right
)


def _safe_text(text: str) -> str:
    """Escape XML-special characters for ReportLab Paragraph."""
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
    )


def _pt_size(run) -> float | None:
    """Return font size in points for a run, or None."""
    if run.font and run.font.size:
        return run.font.size.pt
    return None


def _shape_text_items(shape) -> list[tuple[str, float | None, bool, bool]]:
    """Extract (text, pt_size, bold, is_bullet) tuples from a shape."""
    items: list[tuple[str, float | None, bool, bool]] = []
    if not shape.has_text_frame:
        return items
    for para in shape.text_frame.paragraphs:
        raw = para.text.strip()
        if not raw:
            continue
        # detect bullet by paragraph level or presence of <a:buChar> / <a:buAutoNum>
        pPr = para._p.pPr
        has_bu_elem = False
        if pPr is not None:
            ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
            has_bu_elem = (
                pPr.find(f"{{{ns}}}buChar") is not None
                or pPr.find(f"{{{ns}}}buAutoNum") is not None
            )
        is_bullet = (para.level > 0) or has_bu_elem
        # aggregate font properties from runs
        sizes = [_pt_size(r) for r in para.runs if _pt_size(r)]
        pt = sizes[0] if sizes else None
        bolds = [r.font.bold for r in para.runs if r.font.bold is not None]
        bold = bool(bolds and bolds[0])
        items.append((raw, pt, bold, is_bullet))
    return items


def _classify_shape(shape, slide_width_emu: int) -> str:
    """Classify a shape as 'title', 'subtitle', or 'body'."""
    # shapes that are positioned in the top ~20 % of the slide
    try:
        top_ratio = shape.top / slide_width_emu
        if top_ratio < 0.15:
            return "title"
        if top_ratio < 0.30:
            return "subtitle"
    except Exception:
        pass
    try:
        ph = shape.placeholder_format
        if ph is not None:
            idx = ph.idx
            if idx == 0:
                return "title"
            if idx == 1:
                return "subtitle"
    except Exception:
        pass
    return "body"


# ---------------------------------------------------------------------------
# Pure-Python conversion
# ---------------------------------------------------------------------------

def convert_pptx_to_pdf(input_path: Path, output_path: Path) -> None:
    """Convert a PPTX file to a readable PDF using python-pptx + reportlab."""
    prs = Presentation(str(input_path))
    slide_width_emu = prs.slide_width

    doc = SimpleDocTemplate(
        str(output_path),
        pagesize=landscape(A4),
        leftMargin=2 * cm,
        rightMargin=2 * cm,
        topMargin=2 * cm,
        bottomMargin=1.5 * cm,
    )

    story = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_elements: list = []

        # --- slide number ---
        slide_elements.append(
            Paragraph(f"Slide {slide_idx}", _STYLE_SLIDE_NUM)
        )
        slide_elements.append(Spacer(1, 4 * mm))

        # --- collect all text from all shapes ---
        title_paras: list[Paragraph] = []
        subtitle_paras: list[Paragraph] = []
        body_paras: list[Paragraph] = []

        for shape in slide.shapes:
            role = _classify_shape(shape, slide_width_emu)
            items = _shape_text_items(shape)
            for text, pt, bold, is_bullet in items:
                safe = _safe_text(text)
                if role == "title":
                    title_paras.append(Paragraph(safe, _STYLE_TITLE))
                elif role == "subtitle":
                    subtitle_paras.append(Paragraph(safe, _STYLE_SUBTITLE))
                else:
                    if is_bullet:
                        bullet_text = f"• {safe}"
                        body_paras.append(
                            Paragraph(bullet_text, _STYLE_BULLET)
                        )
                    elif bold and pt and pt >= 14:
                        body_paras.append(
                            Paragraph(f"<b>{safe}</b>", _STYLE_SUBTITLE)
                        )
                    else:
                        body_paras.append(Paragraph(safe, _STYLE_BODY))

        if not title_paras and not subtitle_paras and not body_paras:
            # blank slide placeholder
            body_paras.append(
                Paragraph("[blank slide]", _STYLE_BODY)
            )

        slide_elements.extend(title_paras)
        slide_elements.extend(subtitle_paras)
        if subtitle_paras or title_paras:
            slide_elements.append(Spacer(1, 6 * mm))
        slide_elements.extend(body_paras)

        story.extend(slide_elements)
        story.append(PageBreak())

    doc.build(story)


# ---------------------------------------------------------------------------
# LibreOffice conversion
# ---------------------------------------------------------------------------

def _find_libreoffice() -> str | None:
    """Return path to soffice/libreoffice executable, or None."""
    for candidate in ("libreoffice", "soffice", "/usr/bin/libreoffice",
                      "/usr/bin/soffice", "/opt/libreoffice/program/soffice"):
        try:
            result = subprocess.run(
                [candidate, "--version"],
                capture_output=True, timeout=10
            )
            if result.returncode == 0:
                return candidate
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue
    return None


def convert_with_libreoffice(input_path: Path, output_path: Path) -> None:
    """Convert using LibreOffice headless mode."""
    lo = _find_libreoffice()
    if lo is None:
        raise RuntimeError(
            "LibreOffice is not installed or not found in PATH. "
            "Install it with: sudo apt install libreoffice"
        )
    out_dir = output_path.parent
    cmd = [
        lo,
        "--headless",
        "--convert-to", "pdf",
        "--outdir", str(out_dir),
        str(input_path),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=120)
    if result.returncode != 0:
        raise RuntimeError(
            f"LibreOffice conversion failed:\n{result.stderr}"
        )
    # LibreOffice names the output after the input stem
    lo_output = out_dir / (input_path.stem + ".pdf")
    if lo_output != output_path:
        lo_output.rename(output_path)


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        prog="slide2pdf",
        description=(
            "Convert a presentation (PPTX/ODP) to a readable PDF.\n\n"
            "By default a pure-Python pipeline is used (python-pptx + "
            "reportlab). Pass --libreoffice to use LibreOffice headless."
        ),
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument("input", help="Input presentation file (.pptx, .odp)")
    parser.add_argument(
        "output",
        nargs="?",
        help="Output PDF file path (default: <input_stem>.pdf in same directory)",
    )
    parser.add_argument(
        "--libreoffice",
        action="store_true",
        default=False,
        help="Use LibreOffice headless as conversion backend",
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    parser = _build_parser()
    args = parser.parse_args(argv)

    input_path = Path(args.input).expanduser().resolve()
    if not input_path.exists():
        print(f"Error: input file not found: {input_path}", file=sys.stderr)
        return 1

    suffix = input_path.suffix.lower()
    if suffix not in {".pptx", ".odp", ".ppt"}:
        print(
            f"Warning: unrecognized extension '{suffix}'. "
            "Proceeding anyway.",
            file=sys.stderr,
        )

    if args.output:
        output_path = Path(args.output).expanduser().resolve()
    else:
        output_path = input_path.with_suffix(".pdf")

    output_path.parent.mkdir(parents=True, exist_ok=True)

    try:
        if args.libreoffice:
            print(f"Converting '{input_path.name}' with LibreOffice…")
            convert_with_libreoffice(input_path, output_path)
        else:
            print(f"Converting '{input_path.name}' with python-pptx + reportlab…")
            convert_pptx_to_pdf(input_path, output_path)
        print(f"PDF written to: {output_path}")
        return 0
    except Exception as exc:  # noqa: BLE001
        print(f"Error: {exc}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    sys.exit(main())
