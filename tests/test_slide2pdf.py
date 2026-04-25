"""Tests for slide2pdf.py"""

import io
import os
import sys
from pathlib import Path

import pytest

# Ensure repo root is on path so we can import slide2pdf
sys.path.insert(0, str(Path(__file__).parent.parent))

import slide2pdf  # noqa: E402


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture()
def simple_pptx(tmp_path):
    """Create a minimal PPTX with two slides and return its path."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide_layout = prs.slide_layouts[0]  # title slide

    # Slide 1 – title + body
    slide1 = prs.slides.add_slide(slide_layout)
    slide1.shapes.title.text = "Hello, Slide2PDF"
    slide1.placeholders[1].text = "A subtitle on slide one"

    # Slide 2 – blank layout with a text box
    blank_layout = prs.slide_layouts[6]
    slide2 = prs.slides.add_slide(blank_layout)
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(3))
    tf = txBox.text_frame
    tf.text = "Bullet point one"
    tf.add_paragraph().text = "Bullet point two"

    pptx_path = tmp_path / "test_presentation.pptx"
    prs.save(str(pptx_path))
    return pptx_path


# ---------------------------------------------------------------------------
# Unit tests – helper functions
# ---------------------------------------------------------------------------

class TestSafeText:
    def test_ampersand(self):
        assert slide2pdf._safe_text("A & B") == "A &amp; B"

    def test_less_than(self):
        assert slide2pdf._safe_text("x < y") == "x &lt; y"

    def test_greater_than(self):
        assert slide2pdf._safe_text("x > y") == "x &gt; y"

    def test_no_special(self):
        assert slide2pdf._safe_text("Hello World") == "Hello World"

    def test_combined(self):
        assert slide2pdf._safe_text("<b>A & B</b>") == "&lt;b&gt;A &amp; B&lt;/b&gt;"


class TestFindLibreOffice:
    def test_returns_none_when_missing(self, monkeypatch):
        """When no LibreOffice binary is present, returns None."""
        import subprocess

        def fake_run(cmd, **kwargs):
            raise FileNotFoundError

        monkeypatch.setattr(subprocess, "run", fake_run)
        assert slide2pdf._find_libreoffice() is None


# ---------------------------------------------------------------------------
# Integration tests – actual conversion
# ---------------------------------------------------------------------------

class TestConvertPptxToPdf:
    def test_output_file_created(self, simple_pptx, tmp_path):
        output = tmp_path / "output.pdf"
        slide2pdf.convert_pptx_to_pdf(simple_pptx, output)
        assert output.exists()

    def test_output_is_valid_pdf(self, simple_pptx, tmp_path):
        output = tmp_path / "output.pdf"
        slide2pdf.convert_pptx_to_pdf(simple_pptx, output)
        # PDF files start with %PDF
        header = output.read_bytes()[:4]
        assert header == b"%PDF"

    def test_output_non_empty(self, simple_pptx, tmp_path):
        output = tmp_path / "output.pdf"
        slide2pdf.convert_pptx_to_pdf(simple_pptx, output)
        assert output.stat().st_size > 1024  # at least 1 KB


# ---------------------------------------------------------------------------
# CLI tests
# ---------------------------------------------------------------------------

class TestCLI:
    def test_missing_input_returns_error(self):
        rc = slide2pdf.main(["nonexistent_file.pptx"])
        assert rc == 1

    def test_successful_conversion(self, simple_pptx, tmp_path):
        output = tmp_path / "cli_out.pdf"
        rc = slide2pdf.main([str(simple_pptx), str(output)])
        assert rc == 0
        assert output.exists()

    def test_default_output_path(self, simple_pptx):
        # When no output is given, file should be next to input
        expected = simple_pptx.with_suffix(".pdf")
        try:
            rc = slide2pdf.main([str(simple_pptx)])
            assert rc == 0
            assert expected.exists()
        finally:
            if expected.exists():
                expected.unlink()

    def test_libreoffice_missing_returns_error(self, simple_pptx, tmp_path, monkeypatch):
        """When LibreOffice is absent, --libreoffice flag should fail gracefully."""
        monkeypatch.setattr(slide2pdf, "_find_libreoffice", lambda: None)
        output = tmp_path / "lo_out.pdf"
        rc = slide2pdf.main(["--libreoffice", str(simple_pptx), str(output)])
        assert rc == 1
